from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── COLORES BUPA ──────────────────────────────────────────────
BLUE        = RGBColor(0x00, 0x66, 0xCC)
BLUE_DARK   = RGBColor(0x00, 0x4E, 0xA3)
BLUE_LIGHT  = RGBColor(0xEB, 0xF4, 0xFF)
BG          = RGBColor(0xF3, 0xF7, 0xFB)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT        = RGBColor(0x0F, 0x21, 0x37)
TEXT2       = RGBColor(0x3D, 0x5A, 0x73)
TEXT3       = RGBColor(0x7A, 0x98, 0xB0)
BORDER      = RGBColor(0xD8, 0xE8, 0xF4)
GRAY_LIGHT  = RGBColor(0xF8, 0xFA, 0xFC)

C_CIMA  = RGBColor(0xBF, 0x26, 0x00)
C_E2E   = RGBColor(0x94, 0x62, 0x00)
C_INT   = RGBColor(0x52, 0x43, 0xAA)
C_BASE  = RGBColor(0x00, 0x87, 0x5A)

GREEN      = RGBColor(0x00, 0x87, 0x5A)
GREEN_BG   = RGBColor(0xE6, 0xF5, 0xEF)
YELLOW     = RGBColor(0x94, 0x62, 0x00)
YELLOW_BG  = RGBColor(0xFF, 0xF3, 0xCD)
RED        = RGBColor(0xBF, 0x26, 0x00)
RED_BG     = RGBColor(0xFF, 0xEB, 0xE6)
PURPLE     = RGBColor(0x52, 0x43, 0xAA)
PURPLE_BG  = RGBColor(0xEA, 0xE6, 0xFF)


def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.width = lw
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def tb(slide, x, y, w, h, text, size=Pt(12), bold=False,
       color=TEXT, align=PP_ALIGN.LEFT, italic=False):
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return t


def footer(slide, w, h):
    rect(slide, 0, h - Inches(0.38), w, Inches(0.38), fill=BLUE_DARK)
    tb(slide, Inches(0.4), h - Inches(0.33), w - Inches(0.8), Inches(0.28),
       "BUPA Chile  ·  QA Lead: Jaime Quiñelen Villar  ·  Pirámide de Testing  ·  Mayo 2026",
       size=Pt(10), color=WHITE)


def header_band(slide, w):
    rect(slide, 0, 0, w, Inches(0.08), fill=BLUE)


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
SH = prs.slide_height
SW = prs.slide_width
blank = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, SW, SH, fill=BG)
header_band(s1, SW)

cx, cy, cw, ch = Inches(2.2), Inches(1.5), Inches(9.0), Inches(4.6)
rect(s1, cx, cy, cw, ch, fill=WHITE, line=BORDER, lw=Pt(1))
rect(s1, cx, cy, Inches(0.07), ch, fill=BLUE)

tb(s1, cx + Inches(0.3), cy + Inches(0.35), cw - Inches(0.5), Inches(0.9),
   "Pirámide de Testing — BUPA Chile",
   size=Pt(32), bold=True, color=TEXT)

tb(s1, cx + Inches(0.3), cy + Inches(1.3), cw - Inches(0.5), Inches(0.4),
   "Modelo Cohn  ·  11 Tipos de Prueba  ·  Stack QA",
   size=Pt(16), color=TEXT2)

rect(s1, cx + Inches(0.3), cy + Inches(1.85), cw - Inches(0.6), Pt(1.2), fill=BORDER)

meta_y = cy + Inches(2.1)
metas = [
    ("Autor",        "Jaime Quiñelen Villar — QA Lead"),
    ("Organización", "BUPA Chile"),
    ("Modelo",       "Pirámide de Testing — Mike Cohn"),
    ("Versión",      "v1.0"),
    ("Fecha",        "Mayo 2026"),
]
for label, val in metas:
    tb(s1, cx + Inches(0.3), meta_y, Inches(1.8), Inches(0.36),
       label + ":", size=Pt(12), bold=True, color=TEXT3)
    tb(s1, cx + Inches(2.2), meta_y, cw - Inches(2.5), Inches(0.36),
       val, size=Pt(12), color=TEXT2)
    meta_y += Inches(0.38)

footer(s1, SW, SH)


# ══════════════════════════════════════════════════════════
# SLIDE 2 — PIRÁMIDE VISUAL
# ══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
rect(s2, 0, 0, SW, SH, fill=BG)
header_band(s2, SW)

tb(s2, Inches(0.5), Inches(0.18), Inches(6), Inches(0.28),
   "ESTRATEGIA DE TESTING", size=Pt(12), bold=True, color=BLUE)
tb(s2, Inches(0.5), Inches(0.46), Inches(9), Inches(0.55),
   "Pirámide de Testing — Modelo Cohn", size=Pt(24), bold=True, color=TEXT)
tb(s2, Inches(0.5), Inches(1.02), Inches(9), Inches(0.32),
   "Mientras más arriba: más caro, más lento y menos casos se justifican.",
   size=Pt(12), color=TEXT2, italic=True)

# ── PIRÁMIDE ──
PY_X   = Inches(0.9)
PY_TOP = Inches(1.45)
PY_W   = Inches(7.8)
GAP    = Inches(0.05)

tiers = [
    # name, sub, tools, types, envs, color, h, wfrac
    ("CIMA — Criterio Humano",       "",
     "Sin herramienta fija  ·  Observación directa",
     "Manual  ·  Exploratory  ·  Usability",
     "UAT", C_CIMA, Inches(1.22), 0.30),

    ("E2E / UI  —  10%",             "",
     "Playwright · Cypress · Appium · Selenium · BrowserStack",
     "Regression  ·  Accessibility  ·  Security",
     "UAT · PROD", C_E2E, Inches(1.35), 0.54),

    ("INTEGRACIÓN / SERVICE  —  20%","",
     "Postman · Newman · n8n · APIs REST",
     "Integration  ·  Performance  ·  Smoke",
     "DEV · UAT · PROD", C_INT, Inches(1.45), 0.77),

    ("BASE  —  70%",                 "",
     "Go test · Jest · Karma · Cypress · Playwright · Azure DevOps CI",
     "Unit Testing  ·  Automated Testing",
     "DEV · UAT · PROD", C_BASE, Inches(1.6), 1.00),
]

cy2 = PY_TOP
for (name, _, tools, types, envs, color, h, wfrac) in tiers:
    w = Inches(PY_W.inches * wfrac)
    x = PY_X + (PY_W - w) / 2
    rect(s2, x, cy2, w, h, fill=color)

    tb(s2, x + Inches(0.15), cy2 + Inches(0.08), w - Inches(0.3), Inches(0.35),
       name, size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tb(s2, x + Inches(0.1), cy2 + Inches(0.44), w - Inches(0.2), Inches(0.28),
       tools, size=Pt(12), color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    tb(s2, x + Inches(0.1), cy2 + Inches(0.72), w - Inches(0.2), Inches(0.28),
       types, size=Pt(12), bold=True,
       color=RGBColor(0xFF, 0xFF, 0xAA), align=PP_ALIGN.CENTER)

    tb(s2, x + Inches(0.1), cy2 + h - Inches(0.38), w - Inches(0.2), Inches(0.3),
       f"[ {envs} ]", size=Pt(12), bold=True,
       color=RGBColor(0x7F, 0xFF, 0xD4), align=PP_ALIGN.CENTER)

    cy2 += h + GAP

# Etiquetas eje
tb(s2, Inches(0.05), PY_TOP, Inches(0.72), Inches(0.5),
   "LENTO\nCARO", size=Pt(12), bold=True, color=RED, align=PP_ALIGN.CENTER)
tb(s2, Inches(0.05), cy2 - Inches(0.55), Inches(0.72), Inches(0.5),
   "RÁPIDO\nBARATO", size=Pt(12), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ── PANEL DERECHO ──
RX = Inches(9.0)
RY = PY_TOP
RW = Inches(4.0)

right_cards = [
    ("CIMA",          "No automatizable",  ["Manual", "Exploratory", "Usability"],     RED,    RED_BG),
    ("E2E · 10%",     "Semi-automatizado", ["Regression", "Accessibility", "Security"], YELLOW, YELLOW_BG),
    ("SERVICE · 20%", "Automatizable",     ["Integration", "Performance", "Smoke"],    PURPLE, PURPLE_BG),
    ("BASE · 70%",    "100% automatizado", ["Unit Testing", "Automated"],               GREEN,  GREEN_BG),
]

card_h = Inches(1.52)
for (badge, lbl, chips, accent, bg) in right_cards:
    rect(s2, RX, RY, RW, card_h, fill=WHITE, line=BORDER, lw=Pt(0.75))
    rect(s2, RX, RY, Inches(0.06), card_h, fill=accent)

    bw = Inches(len(badge) * 0.105 + 0.25)
    rect(s2, RX + Inches(0.18), RY + Inches(0.12), bw, Inches(0.32), fill=bg)
    tb(s2, RX + Inches(0.18), RY + Inches(0.12), bw, Inches(0.32),
       badge, size=Pt(12), bold=True, color=accent, align=PP_ALIGN.CENTER)

    tb(s2, RX + bw + Inches(0.3), RY + Inches(0.12), RW - bw - Inches(0.5), Inches(0.32),
       lbl, size=Pt(12), color=TEXT3, bold=True)

    chip_x = RX + Inches(0.18)
    chip_y = RY + Inches(0.58)
    for chip in chips:
        chip_w = Inches(len(chip) * 0.1 + 0.3)
        rect(s2, chip_x, chip_y, chip_w, Inches(0.32), fill=BG, line=BORDER, lw=Pt(0.5))
        tb(s2, chip_x + Inches(0.06), chip_y, chip_w, Inches(0.32),
           chip, size=Pt(12), color=TEXT2)
        chip_x += chip_w + Inches(0.1)

    RY += card_h + Inches(0.08)

footer(s2, SW, SH)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — MÉTRICAS + EJE
# ══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
rect(s3, 0, 0, SW, SH, fill=BG)
header_band(s3, SW)

tb(s3, Inches(0.5), Inches(0.18), Inches(6), Inches(0.28),
   "MÉTRICAS POR NIVEL", size=Pt(12), bold=True, color=BLUE)
tb(s3, Inches(0.5), Inches(0.46), Inches(9), Inches(0.55),
   "Distribución por Nivel — Pirámide de Testing", size=Pt(24), bold=True, color=TEXT)

kpis = [
    ("Base — Unitarias",      "70%", "Go test · Jest · Karma\nAzure DevOps CI\nCobertura mín. ≥ 80%", C_BASE),
    ("Integración / Service", "20%", "Postman · Newman · n8n\nAPIs REST · Contratos SDD",              C_INT),
    ("E2E / UI",              "10%", "Playwright · Cypress · Appium\nSelenium · BrowserStack",          C_E2E),
    ("Tipos de Testing",      "11",  "Tipos cubiertos\nen el plan QA BUPA",                             C_CIMA),
]

kx, ky, kw, kh = Inches(0.4), Inches(1.2), Inches(3.1), Inches(2.15)
for (name, val, desc, color) in kpis:
    rect(s3, kx, ky, kw, kh, fill=WHITE, line=BORDER, lw=Pt(0.75))
    rect(s3, kx, ky, kw, Inches(0.06), fill=color)
    tb(s3, kx + Inches(0.18), ky + Inches(0.18), kw - Inches(0.3), Inches(0.32),
       name.upper(), size=Pt(12), bold=True, color=TEXT3)
    tb(s3, kx + Inches(0.18), ky + Inches(0.54), kw - Inches(0.3), Inches(0.75),
       val, size=Pt(42), bold=True, color=color, align=PP_ALIGN.LEFT)
    tb(s3, kx + Inches(0.18), ky + Inches(1.4), kw - Inches(0.3), Inches(0.65),
       desc, size=Pt(12), color=TEXT2)
    kx += kw + Inches(0.22)

# ── Tabla eje ──
tby = Inches(3.55)
tb(s3, Inches(0.5), tby, Inches(6), Inches(0.38),
   "Eje de la Pirámide", size=Pt(18), bold=True, color=TEXT)

headers3 = ["Atributo", "Cima", "E2E", "Service", "Base"]
rows3 = [
    ["Velocidad",         "Lento",  "Medio",      "Medio",         "Rápido"],
    ["Costo",             "Alto",   "Medio-Alto", "Medio",         "Bajo"],
    ["Cantidad de casos", "Pocos",  "Moderado",   "Moderado",      "Muchos"],
    ["Automatización",    "No",     "Parcial",    "Parcial-Total", "Total"],
    ["Feedback",          "Tardío", "Tardío",     "Rápido",        "Inmediato"],
]
col_w3 = [Inches(2.5), Inches(1.9), Inches(1.9), Inches(2.2), Inches(1.9)]
tx3, ty3, rh3 = Inches(0.5), tby + Inches(0.45), Inches(0.46)

hx3 = tx3
for i, h in enumerate(headers3):
    rect(s3, hx3, ty3, col_w3[i], rh3, fill=BLUE)
    tb(s3, hx3 + Inches(0.1), ty3 + Inches(0.08), col_w3[i], rh3,
       h.upper(), size=Pt(12), bold=True, color=WHITE)
    hx3 += col_w3[i]
ty3 += rh3

accents3 = [None, C_CIMA, C_E2E, C_INT, C_BASE]
for row in rows3:
    rx3 = tx3
    for i, cell in enumerate(row):
        bg3 = GRAY_LIGHT if i == 0 else WHITE
        rect(s3, rx3, ty3, col_w3[i], rh3, fill=bg3, line=BORDER, lw=Pt(0.5))
        fc3 = accents3[i] if (i > 0 and accents3[i]) else TEXT2
        tb(s3, rx3 + Inches(0.1), ty3 + Inches(0.08), col_w3[i] - Inches(0.12), rh3,
           cell, size=Pt(12), color=fc3, bold=(i == 0))
        rx3 += col_w3[i]
    ty3 += rh3

footer(s3, SW, SH)


# ══════════════════════════════════════════════════════════
# SLIDE 4 — MAPA COMPLETO 11 TIPOS
# ══════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
rect(s4, 0, 0, SW, SH, fill=BG)
header_band(s4, SW)

tb(s4, Inches(0.4), Inches(0.18), Inches(6), Inches(0.28),
   "DETALLE COMPLETO", size=Pt(12), bold=True, color=BLUE)
tb(s4, Inches(0.4), Inches(0.46), Inches(10), Inches(0.52),
   "Mapa Completo — 11 Tipos en la Pirámide", size=Pt(24), bold=True, color=TEXT)

table_data = [
    ("Unit Testing",         "Base",    "Go test · Jest · Karma",      "100% Auto", "DEV",          "Dev Team",         C_BASE,  GREEN_BG),
    ("Automated Testing",    "Base",    "Cypress · Playwright",         "100% Auto", "UAT · PROD",   "QA Lead + Dev",    C_BASE,  GREEN_BG),
    ("Smoke Testing",        "Service", "Cypress (subset crítico)",     "100% Auto", "UAT · PROD",   "QA Lead",          C_INT,   PURPLE_BG),
    ("Integration Testing",  "Service", "Postman · Newman · n8n",       "Parcial",   "DEV·UAT·PROD", "QA Lead + DevOps", C_INT,   PURPLE_BG),
    ("Performance Testing",  "Service", "Lighthouse · k6 · Artillery",  "Parcial",   "UAT",          "QA Lead + DevOps", C_INT,   PURPLE_BG),
    ("Regression Testing",   "E2E",     "Cypress · TestRail",           "Parcial",   "UAT · PROD",   "QA Lead",          C_E2E,   YELLOW_BG),
    ("Accessibility Testing","E2E",     "cypress-axe · NVDA",           "Parcial",   "UAT",          "QA Lead",          C_E2E,   YELLOW_BG),
    ("Security Testing",     "E2E",     "curl · openssl · OWASP ZAP",  "Parcial",   "UAT · PROD",   "QA Lead + Dev",    C_E2E,   YELLOW_BG),
    ("Manual Testing",       "Cima",    "Navegador · DevTools",         "Manual",    "UAT · PROD",   "QA Lead",          C_CIMA,  RED_BG),
    ("Exploratory Testing",  "Cima",    "Charter de exploración",       "Manual",    "UAT",          "QA Lead",          C_CIMA,  RED_BG),
    ("Usability Testing",    "Cima",    "Usuarios",                     "Manual",    "UAT",          "QA Lead",          C_CIMA,  RED_BG),
]

col_w4  = [Inches(2.2), Inches(0.95), Inches(2.5), Inches(1.15), Inches(1.3), Inches(1.55)]
headers4 = ["Tipo de Testing", "Nivel", "Herramienta", "Automatización", "Ambiente", "Responsable"]
tx4, ty4, rh4 = Inches(0.28), Inches(1.1), Inches(0.48)

hx4 = tx4
for i, h in enumerate(headers4):
    rect(s4, hx4, ty4, col_w4[i], rh4, fill=BLUE)
    tb(s4, hx4 + Inches(0.07), ty4 + Inches(0.1), col_w4[i] - Inches(0.08), rh4,
       h.upper(), size=Pt(12), bold=True, color=WHITE)
    hx4 += col_w4[i]
ty4 += rh4

auto_fg = {"100% Auto": GREEN,  "Parcial": BLUE,       "Manual": YELLOW}
auto_bg = {"100% Auto": GREEN_BG, "Parcial": BLUE_LIGHT, "Manual": YELLOW_BG}

for i, (tipo, nivel, herr, auto, amb, resp, accent, _) in enumerate(table_data):
    row_bg = WHITE if i % 2 == 0 else GRAY_LIGHT
    rx4 = tx4
    for j, cell in enumerate([tipo, nivel, herr, auto, amb, resp]):
        rect(s4, rx4, ty4, col_w4[j], rh4, fill=row_bg, line=BORDER, lw=Pt(0.4))
        if j == 0:
            rect(s4, rx4 + Inches(0.08), ty4 + Inches(0.19), Inches(0.1), Inches(0.1), fill=accent)
            tb(s4, rx4 + Inches(0.24), ty4 + Inches(0.09), col_w4[j] - Inches(0.28), rh4,
               cell, size=Pt(12), bold=True, color=TEXT)
        elif j == 1:
            tb(s4, rx4 + Inches(0.08), ty4 + Inches(0.09), col_w4[j] - Inches(0.1), rh4,
               cell, size=Pt(12), bold=True, color=accent)
        elif j == 3:
            abg = auto_bg.get(cell, BG)
            afg = auto_fg.get(cell, TEXT2)
            rect(s4, rx4 + Inches(0.06), ty4 + Inches(0.1), Inches(0.98), Inches(0.28), fill=abg, line=BORDER, lw=Pt(0.4))
            tb(s4, rx4 + Inches(0.06), ty4 + Inches(0.1), Inches(0.98), Inches(0.28),
               cell, size=Pt(12), bold=True, color=afg, align=PP_ALIGN.CENTER)
        else:
            tb(s4, rx4 + Inches(0.08), ty4 + Inches(0.09), col_w4[j] - Inches(0.1), rh4,
               cell, size=Pt(12), color=TEXT2)
        rx4 += col_w4[j]
    ty4 += rh4

footer(s4, SW, SH)


# ══════════════════════════════════════════════════════════
# GUARDAR
# ══════════════════════════════════════════════════════════
out = r"c:\Quality_Assurance_IA\QA-Engineer\estandarizacion\piramide-testing.pptx"
prs.save(out)
print(f"PPTX generado: {out}")
